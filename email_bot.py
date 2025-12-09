import os
import imaplib
import smtplib
import ssl
import tempfile
import re
from email.message import EmailMessage
from email.parser import BytesParser
from email.policy import default as default_policy
from typing import Callable, List, Optional, Tuple

from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dotenv import load_dotenv

# Load environment before importing modules that rely on OPENAI_API_KEY
load_dotenv()
from agent_runner import run_agent_workflow  # noqa: E402

IMAP_HOST = os.getenv("EMAIL_IMAP_HOST", "imap.gmail.com")
SMTP_HOST = os.getenv("EMAIL_SMTP_HOST", "smtp.gmail.com")
EMAIL_ADDRESS = os.getenv("EMAIL_ADDRESS")
EMAIL_PASSWORD = os.getenv("EMAIL_PASSWORD")


# ---------- Helpers: PPTX → structured JSON ----------

def pptx_to_struct(path: str) -> dict:
    """Convert a .pptx file into {slides: [{index, text}...]}."""
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            texts.extend(_extract_shape_text(shape))

        notes_text = _extract_notes_text(slide)
        if notes_text:
            texts.append(notes_text)

        slides.append({
            "index": i,
            "text": "\n".join(t for t in texts if t).strip()
        })
    return {"slides": slides}


def pdf_to_struct(path: str) -> dict:
    """Convert a PDF file into the same slide structure (page-per-slide)."""
    reader = PdfReader(path)
    slides = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        slides.append({
            "index": i,
            "text": text.strip()
        })
    return {"slides": slides}


def attachment_to_struct(attachment: dict) -> dict:
    filename = attachment.get("filename") or ""
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return pdf_to_struct(attachment["path"])
    return pptx_to_struct(attachment["path"])


def _extract_notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes_frame = slide.notes_slide.notes_text_frame
    if not notes_frame:
        return ""
    text = "\n".join(
        paragraph.text.strip()
        for paragraph in notes_frame.paragraphs
        if paragraph.text and paragraph.text.strip()
    )
    return f"[Notes] {text}" if text else ""


def _extract_shape_text(shape) -> List[str]:
    """
    Recursively collect text from any shape, including grouped overlays,
    tables, and charts. These overlays often carry reviewer comments.
    """
    texts: List[str] = []

    # Group shapes contain nested shapes that may have their own text.
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(_extract_shape_text(child))

    # Standard text frames
    if getattr(shape, "has_text_frame", False):
        frame_text = "\n".join(
            paragraph.text.strip()
            for paragraph in shape.text_frame.paragraphs
            if paragraph.text and paragraph.text.strip()
        )
        if frame_text:
            texts.append(frame_text)
    elif hasattr(shape, "text"):  # Fallback for placeholders without text_frame
        text = (shape.text or "").strip()
        if text:
            texts.append(text)

    # Tables can contain reviewer notes in cells
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    texts.append(cell_text)

    # Chart titles/data labels sometimes hold textual comments
    if getattr(shape, "has_chart", False):
        chart = shape.chart
        if chart.has_title:
            chart_title = chart.chart_title.text_frame.text.strip()
            if chart_title:
                texts.append(chart_title)
        for series in chart.series:
            if not getattr(series, "data_labels", None):
                continue
            for point in getattr(series, "points", []):
                data_label = getattr(point, "data_label", None)
                if not data_label or not getattr(data_label, "has_text_frame", False):
                    continue
                label_text = "\n".join(
                    paragraph.text.strip()
                    for paragraph in data_label.text_frame.paragraphs
                    if paragraph.text and paragraph.text.strip()
                )
                if label_text:
                    texts.append(label_text)

    return texts


# ---------- Helpers: Email parsing ----------

def extract_body_and_attachments(raw_msg: bytes):
    """Return (body_text, [attachment_info]) for a raw email."""
    msg = BytesParser(policy=default_policy).parsebytes(raw_msg)

    # Body text (prefer plain)
    body_text = ""
    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            disp = part.get_content_disposition()
            if ctype == "text/plain" and disp is None:
                body_text = part.get_content()
                break
    else:
        body_text = msg.get_content()

    # Attachments
    attachment_infos = []
    for part in msg.walk():
        disp = part.get_content_disposition()
        if disp == "attachment":
            filename = part.get_filename()
            if not filename:
                continue
            # save to temp file
            data = part.get_payload(decode=True)
            fd, tmp_path = tempfile.mkstemp(suffix=filename)
            with os.fdopen(fd, "wb") as f:
                f.write(data)
            attachment_infos.append({
                "path": tmp_path,
                "filename": filename
            })

    from_addr = msg["From"]
    subject = msg["Subject"]
    message_id = msg["Message-ID"]

    return body_text.strip(), attachment_infos, from_addr, subject, message_id


# ---------- Call your Agent Builder workflow ----------

def run_agent(email_text: str, original_doc: dict, revised_doc: dict) -> dict:
    """
    Call the Agent Builder workflow you exported as `agent_workflow.py`.

    The generated file includes an example of how to run it – usually something like:
       result = workflow.run(input={...})
    or:
       result = runner.run(input={...})

    Adjust this function to match the exact call from that file.
    """
    # ⬇️ EXAMPLE – you must align with the code that Agent Builder generated.
    # Look in agent_workflow.py for the "main" run call and copy it here.

    # For example, if agent_workflow defines an async `run` function:
    #
    #   from agents import run
    #   async def run_workflow(input):
    #       ...
    #
    # You might do something like:
    #
    #   result = asyncio.run(agent_workflow.run_workflow(
    #       {"email_text": email_text,
    #        "original_doc": original_doc,
    #        "revised_doc": revised_doc}
    #   ))
    #
    # And ensure it returns the JSON matching your Node 3 schema.

    # For now, we’ll assume there is a synchronous helper:
    lowered = (email_text or "").lower()
    only_tick = "only tie out" in lowered or "only tick and tie" in lowered
    run_tick_tie = "tie out" in lowered or "tick and tie" in lowered

    result = run_agent_workflow(
        email_text=email_text,
        original_doc=original_doc,
        revised_doc=revised_doc,
        run_tick_tie=run_tick_tie,
    )

    if only_tick:
        return {
            "tags": [],
            "email_comments": [],
            "tick_tie": result.get("tick_tie") or {"ties_out": [], "check": []},
            "only_tick": True,
        }

    result["only_tick"] = False
    return result


# ---------- Formatting the reply email ----------

def format_summary(tags_comments: list, email_comments: list, tick_tie: dict | None = None, show_comments: bool = True) -> str:
    """Turn the agent JSON into a banker-style email body."""
    def buckets(comments: list):
        return (
            _sort_by_slide([c for c in comments if c["status"] == "implemented"]),
            _sort_by_slide([c for c in comments if c["status"] == "partially_implemented"]),
            _sort_by_slide([c for c in comments if c["status"] == "not_implemented"]),
            _sort_by_slide([c for c in comments if c["status"] == "unclear"]),
        )

    tags_impl, tags_part, tags_miss, tags_unclear = buckets(tags_comments)
    email_impl, email_part, email_miss, email_unclear = buckets(email_comments)

    lines = []
    if show_comments:
        lines.append("Coverage summary:")
        lines.append(f"- Tags: {len(tags_impl)} implemented, {len(tags_part)} partial, {len(tags_miss)} not, {len(tags_unclear)} unclear")
        lines.append(f"- Email: {len(email_impl)} implemented, {len(email_part)} partial, {len(email_miss)} not, {len(email_unclear)} unclear\n")

        for title, impl, part, miss, unclear in [
            ("Tags", tags_impl, tags_part, tags_miss, tags_unclear),
            ("Email", email_impl, email_part, email_miss, email_unclear),
        ]:
            lines.append(f"{title}:")
            section = [
                ("Implemented", impl, False),
                ("Partially implemented", part, True),
                ("Not implemented", miss, True),
                ("Unclear / needs human review", unclear, False),
            ]
            for label, bucket, include_suggestion in section:
                if not bucket:
                    continue
                lines.append(f"{label}:")
                lines.extend(_format_bucket_by_slide(bucket, include_suggestion))
                lines.append("")
            lines.append("")

    if tick_tie is not None:
        lines.append("Tick & tie review:")
        if isinstance(tick_tie, list):
            ties_out = tick_tie
            check = []
        else:
            ties_out = tick_tie.get("ties_out") or []
            check = tick_tie.get("check") or []
        if not ties_out and not check:
            lines.append("- No duplicate values in deck.")
            lines.append("")
        else:
            if ties_out:
                lines.append("- Consistent:")
                for item in ties_out:
                    pages = ", ".join(str(p) for p in item.get("pages", []))
                    lines.append(f"  - {item.get('metric_label')}: {item.get('canonical_value')} (pages {pages})")
            if check:
                lines.append("- Inconsistencies to review:")
                for item in check:
                    lines.append(f"  - {item.get('metric_label')}:")
                    for vb in item.get("values_by_page", []):
                        lines.append(f"    - Page {vb.get('page')}: {vb.get('value')}")
                    if item.get("reason"):
                        lines.append(f"    Reason: {item['reason']}")

    return "\n".join(lines)


def _sort_by_slide(comments: list) -> list:
    def slide_key(comment):
        refs = comment.get("slide_refs") or []
        return min(refs) if refs else float("inf")

    return sorted(comments, key=slide_key)


def _format_bucket_by_slide(comments: list, include_suggestion: bool) -> list:
    grouped = {}
    for c in comments:
        refs = c.get("slide_refs") or ["n/a"]
        key = ", ".join(str(r) for r in refs)
        grouped.setdefault(key, []).append(c)

    lines: list[str] = []
    for slide in sorted(grouped, key=lambda s: float(s.split(",")[0]) if s != "n/a" else float("inf")):
        lines.append(f"- Slide {slide}:")
        for c in _sort_by_slide(grouped[slide]):
            lines.append(f"    - Comment {c['id']}: {c['text']}")
            lines.append(f"      Reason: {c['reason']}")
            if include_suggestion and c.get("suggestion"):
                lines.append(f"      Suggestion: {c['suggestion']}")
    return lines


# ---------- Main processing: read email → agent → send email ----------

def process_one_email(raw_msg: bytes):
    email_text, attachments, from_addr, subject, message_id = \
        extract_body_and_attachments(raw_msg)

    valid_attachments = [
        att for att in attachments
        if att["filename"] and os.path.splitext(att["filename"])[1].lower() in {".pptx", ".pdf"}
    ]

    if len(valid_attachments) == 0:
        print("No supported attachments found; skipping.")
        return
    elif len(valid_attachments) == 1:
        revised_att = valid_attachments[0]
        original_doc = {"slides": []}
        revised_doc = attachment_to_struct(revised_att)
    else:
        original_att, revised_att = _choose_original_and_revised(valid_attachments)
        original_doc = attachment_to_struct(original_att)
        revised_doc = attachment_to_struct(revised_att)

    result = run_agent(email_text, original_doc, revised_doc)
    tags_comments = result.get("tags", [])
    email_comments = result.get("email_comments", [])
    tick_tie = result.get("tick_tie")
    only_tick = result.get("only_tick", False)

    summary = format_summary(tags_comments, email_comments, tick_tie, show_comments=not only_tick)

    # Send reply to yourself (or to original sender)
    send_email(
        to_addr=from_addr,
        subject=f"Re: {subject} [Bifocal Review]",
        body=summary,
    )


def send_email(to_addr: str, subject: str, body: str):
    msg = EmailMessage()
    msg["From"] = EMAIL_ADDRESS
    msg["To"] = to_addr
    msg["Subject"] = subject
    msg.set_content(body)

    context = ssl.create_default_context()
    with smtplib.SMTP_SSL(SMTP_HOST, 465, context=context) as server:
        server.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
        server.send_message(msg)
        print(f"Sent coverage email to {to_addr}")


def poll_inbox():
    """Simple one-shot poll: fetch unread emails, process, mark as seen."""
    mail = imaplib.IMAP4_SSL(IMAP_HOST)
    try:
        mail.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
        mail.select("INBOX")

        typ, data = mail.search(None, "UNSEEN")
        if typ != "OK":
            print("No messages found.")
            return

        for num in data[0].split():
            try:
                typ, msg_data = mail.fetch(num, "(RFC822)")
                if typ != "OK":
                    continue
                raw_msg = msg_data[0][1]
                print(f"Processing message UID {num.decode()}")
                process_one_email(raw_msg)
                # Mark as seen (already handled)
                mail.store(num, "+FLAGS", "\\Seen")
            except imaplib.IMAP4.abort as exc:
                # Connection dropped mid-loop; log and break out
                print(f"IMAP connection aborted while processing UID {num.decode()}: {exc}")
                break
    finally:
        try:
            mail.close()
        except Exception:
            pass
        try:
            mail.logout()
        except Exception:
            pass


def _choose_original_and_revised(attachments: List[dict]) -> Tuple[dict, dict]:
    """
    Determine original vs revised PPTX based on version numbers in filenames.
    When two or more attachments include `_vNN`, the lowest number is original
    and the highest is revised. Otherwise fall back to alphabetical order.
    """
    versioned = [
        (att, _parse_version(att.get("filename")))
        for att in attachments
    ]
    valid = [item for item in versioned if item[1] is not None]
    if len(valid) >= 2:
        valid.sort(key=lambda item: item[1])
        return valid[0][0], valid[-1][0]

    sorted_atts = sorted(
        attachments,
        key=lambda att: (att.get("filename") or "").lower()
    )
    return sorted_atts[0], sorted_atts[-1]


def _parse_version(filename: Optional[str]) -> Optional[int]:
    """Extract the integer after '_v' / '-v' etc. Return None if absent."""
    if not filename:
        return None
    match = re.search(r"[._-]v(\d+)", filename, re.IGNORECASE)
    if match:
        return int(match.group(1))
    return None


if __name__ == "__main__":
    poll_inbox()
